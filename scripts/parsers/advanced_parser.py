#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys
import json
import os
from pathlib import Path

# PDF 처리
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# Office 문서 처리
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pandas as pd
except ImportError:
    pd = None

def parse_pdf(file_path):
    """PDF 파일 파싱"""
    if not fitz:
        raise ImportError("PyMuPDF가 설치되지 않았습니다: pip install PyMuPDF")

    doc = fitz.open(file_path)
    result = {
        "content": "",
        "tables": [],
        "images": [],
        "metadata": {
            "pages": doc.page_count,
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
            "subject": doc.metadata.get("subject", "")
        }
    }

    for page_num in range(doc.page_count):
        page = doc[page_num]

        # 텍스트 추출
        page_text = page.get_text()
        result["content"] += f"\n--- 페이지 {page_num + 1} ---\n{page_text}"

        # 표 추출
        try:
            tables = page.find_tables()
            for table_index, table in enumerate(tables):
                table_data = table.extract()
                if table_data:  # 빈 표가 아닌 경우만
                    result["tables"].append({
                        "page": page_num + 1,
                        "table_index": table_index + 1,
                        "data": table_data,
                        "row_count": len(table_data),
                        "col_count": len(table_data[0]) if table_data else 0
                    })
        except Exception as e:
            print(f"표 추출 중 오류 (페이지 {page_num + 1}): {str(e)}", file=sys.stderr)

        # 이미지 정보
        try:
            images = page.get_images()
            for img_index, img in enumerate(images):
                result["images"].append({
                    "page": page_num + 1,
                    "image_index": img_index + 1,
                    "xref": img[0],
                    "bbox": list(img[1:5]) if len(img) >= 5 else None
                })
        except Exception as e:
            print(f"이미지 추출 중 오류 (페이지 {page_num + 1}): {str(e)}", file=sys.stderr)

    doc.close()
    return result

def parse_docx(file_path):
    """DOCX 파일 파싱"""
    if not Document:
        raise ImportError("python-docx가 설치되지 않았습니다: pip install python-docx")

    doc = Document(file_path)
    result = {
        "content": "",
        "tables": [],
        "images": [],
        "metadata": {
            "title": doc.core_properties.title or "",
            "author": doc.core_properties.author or "",
            "subject": doc.core_properties.subject or ""
        }
    }

    # 표 추출
    for table_index, table in enumerate(doc.tables):
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            table_data.append(row_data)

        if table_data and any(any(cell for cell in row) for row in table_data):  # 빈 표가 아닌 경우
            result["tables"].append({
                "table_index": table_index + 1,
                "data": table_data,
                "row_count": len(table_data),
                "col_count": len(table_data[0]) if table_data else 0
            })

    # 텍스트 추출 (표 제외)
    for para in doc.paragraphs:
        if para.text.strip():
            result["content"] += para.text + "\n"

    # 이미지 정보 (관계 정보를 통해)
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                result["images"].append({
                    "relationship_id": rel.rId,
                    "target": rel.target_ref
                })
    except Exception as e:
        print(f"이미지 정보 추출 중 오류: {str(e)}", file=sys.stderr)

    return result

def parse_pptx(file_path):
    """PPTX 파일 파싱"""
    if not Presentation:
        raise ImportError("python-pptx가 설치되지 않았습니다: pip install python-pptx")

    prs = Presentation(file_path)
    result = {
        "content": "",
        "tables": [],
        "images": [],
        "slides": [],
        "metadata": {
            "title": prs.core_properties.title or "",
            "author": prs.core_properties.author or "",
            "slide_count": len(prs.slides)
        }
    }

    for slide_index, slide in enumerate(prs.slides):
        slide_content = f"\n--- 슬라이드 {slide_index + 1} ---\n"
        slide_tables = []
        slide_images = []

        for shape in slide.shapes:
            # 텍스트 추출
            if hasattr(shape, "text") and shape.text.strip():
                slide_content += shape.text + "\n"

            # 표 추출
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)

                if table_data:
                    table_info = {
                        "slide": slide_index + 1,
                        "data": table_data,
                        "row_count": len(table_data),
                        "col_count": len(table_data[0]) if table_data else 0
                    }
                    result["tables"].append(table_info)
                    slide_tables.append(table_info)

            # 이미지 정보
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_info = {
                    "slide": slide_index + 1,
                    "shape_id": shape.shape_id,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
                result["images"].append(image_info)
                slide_images.append(image_info)

        result["content"] += slide_content
        result["slides"].append({
            "slide_number": slide_index + 1,
            "content": slide_content,
            "table_count": len(slide_tables),
            "image_count": len(slide_images)
        })

    return result

def parse_xlsx(file_path):
    """XLSX 파일 파싱"""
    if not pd:
        raise ImportError("pandas가 설치되지 않았습니다: pip install pandas openpyxl")

    try:
        # 모든 시트 읽기
        excel_file = pd.ExcelFile(file_path)
        result = {
            "content": "",
            "tables": [],
            "images": [],
            "metadata": {
                "sheet_names": excel_file.sheet_names,
                "sheet_count": len(excel_file.sheet_names)
            }
        }

        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)

            # 빈 시트가 아닌 경우만 처리
            if not df.empty:
                # 테이블 데이터로 저장
                table_data = [df.columns.tolist()] + df.fillna('').astype(str).values.tolist()

                result["tables"].append({
                    "sheet_name": sheet_name,
                    "data": table_data,
                    "row_count": len(table_data),
                    "col_count": len(table_data[0]) if table_data else 0
                })

                # 텍스트 형태로도 저장
                result["content"] += f"\n--- 시트: {sheet_name} ---\n"
                result["content"] += df.to_string(index=False) + "\n"

        excel_file.close()
        return result

    except Exception as e:
        raise Exception(f"Excel 파일 파싱 실패: {str(e)}")

def main():
    if len(sys.argv) != 2:
        print(json.dumps({"error": "사용법: python advanced_parser.py <파일경로>"}, ensure_ascii=False))
        sys.exit(1)

    file_path = sys.argv[1]

    if not os.path.exists(file_path):
        print(json.dumps({"error": f"파일이 존재하지 않습니다: {file_path}"}, ensure_ascii=False))
        sys.exit(1)

    file_ext = Path(file_path).suffix.lower()

    try:
        if file_ext == '.pdf':
            result = parse_pdf(file_path)
        elif file_ext == '.docx':
            result = parse_docx(file_path)
        elif file_ext == '.pptx':
            result = parse_pptx(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            result = parse_xlsx(file_path)
        else:
            result = {"error": f"지원하지 않는 파일 형식: {file_ext}"}

        # 성공적으로 파싱된 경우 통계 정보 추가
        if "error" not in result:
            result["parsing_stats"] = {
                "content_length": len(result.get("content", "")),
                "table_count": len(result.get("tables", [])),
                "image_count": len(result.get("images", [])),
                "word_count": len(result.get("content", "").split()) if result.get("content") else 0
            }

        print(json.dumps(result, ensure_ascii=False, indent=2))

    except Exception as e:
        error_result = {
            "error": f"{type(e).__name__}: {str(e)}",
            "file_path": file_path,
            "file_type": file_ext
        }
        print(json.dumps(error_result, ensure_ascii=False))
        sys.exit(1)

if __name__ == "__main__":
    main()